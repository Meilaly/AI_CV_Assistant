# 导入所需的库
from transformers import AutoTokenizer, AutoModelForCausalLM
import torch
import streamlit as st
from peft import PeftModel
import json
import pandas as pd
import docx  # 用于处理.docx文件
from pdfplumber import open as pdf_open  # 用于处理.pdf文件
from pptx import Presentation  # 用于处理.pptx文件

# 源大模型下载
from modelscope import snapshot_download
model_dir = snapshot_download('IEITYuan/Yuan2-2B-Mars-hf', cache_dir='./')

# 定义模型路径
path = './IEITYuan/Yuan2-2B-Mars-hf'
lora_path = './output/Yuan2.0-2B_lora_bf16/checkpoint-51'

# 定义模型数据类型
torch_dtype = torch.bfloat16 # A10
# torch_dtype = torch.float16 # P100

# 定义一个函数，用于获取模型和tokenizer
@st.cache_resource
def get_model():
    print("Creat tokenizer...")
    tokenizer = AutoTokenizer.from_pretrained(path, add_eos_token=False, add_bos_token=False, eos_token='<eod>')
    tokenizer.add_tokens(['<sep>', '<pad>', '<mask>', '<predict>', '<FIM_SUFFIX>', '<FIM_PREFIX>', '<FIM_MIDDLE>',
                          '<commit_before>', '<commit_msg>', '<commit_after>', '<jupyter_start>', '<jupyter_text>',
                          '<jupyter_code>', '<jupyter_output>', '<empty_output>'], special_tokens=True)

    print("Creat model...")
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    model = AutoModelForCausalLM.from_pretrained(
        path,
        torch_dtype=torch.bfloat16 if device.type == 'cuda' else torch.float32,
        trust_remote_code=True
    )
    model.to(device)  # 移动模型到指定的设备
    model = PeftModel.from_pretrained(model, model_id=lora_path)

    return tokenizer, model

# 加载model和tokenizer
tokenizer, model = get_model()


template = '''
# 任务描述
假设你是一个AI简历助手，能从简历中识别出所有的命名实体，并以json格式返回结果。

# 任务要求
实体的类别包括：姓名、国籍、种族、职位、教育背景、专业、组织名、地名。
返回的json格式是一个字典，其中每个键是实体的类别，值是一个列表，包含实体的文本。

# 样例
输入：
张三，男，中国籍，工程师
输出：
{"姓名": ["张三"], "国籍": ["中国"], "职位": ["工程师"]}

# 当前简历
query

# 任务重述
请参考样例，按照任务要求，识别出当前简历中所有的命名实体，并以json格式返回结果。
'''
# 创建一个标题和一个副标题
st.title("💬 Yuan2.0 AI简历助手")

# 在聊天界面上显示模型的输出
st.write("请上传简历文件：")

# 如果用户上传了文件，则执行以下操作
uploaded_file = st.file_uploader("上传文件", type=["pdf", "docx", "pptx", "txt"])

# 在聊天界面上显示模型的输出
st.chat_message("assistant").write(f"请输入简历文本：")

# 读取文件的函数
def read_file(file):
    if file.type == 'application/pdf':
        with pdf_open(file) as pdf:
            text = '\n'.join(page.extract_text() for page in pdf.pages)
    elif file.type.startswith('application/vnd.openxmlformats-officedocument'):
        if file.name.endswith('.docx'):
            doc = docx.Document(file)
            text = '\n'.join([p.text for p in doc.paragraphs])
        elif file.name.endswith('.pptx'):
            prs = Presentation(file)
            text = ''
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + '\n'
    elif file.type == 'text/plain':
        text = file.read().decode('utf-8')
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, DOCX, PPTX, or TXT file.")

    return text

# 如果用户在聊天输入框中输入了内容，则执行以下操作
if query := st.chat_input():
    # 在聊天界面上显示用户的输入
    st.chat_message("user").write(query)

    # 使用st.empty()创建一个可更新的区域
    empty_placeholder = st.empty()

    # 在该区域内显示“正在提取简历信息，请稍候...”
    with empty_placeholder.container():
        st.write("正在提取简历信息，请稍候...")

    # 调用模型
    prompt = template.replace('query', query).strip()
    prompt += "<sep>"
    inputs = tokenizer(prompt, return_tensors="pt")["input_ids"].cuda()
    outputs = model.generate(inputs, do_sample=False, max_length=1024)  # 设置解码方式和最大生成长度
    output = tokenizer.decode(outputs[0])
    response = output.split("<sep>")[-1].replace("<eod>", '').strip()

    # 清空“正在提取简历信息，请稍候...”的信息
    empty_placeholder.empty()

    # 创建一个新的st.chat_message并显示解析后的表格
    new_placeholder = st.chat_message("assistant")
    with new_placeholder.container():
        df = pd.DataFrame(json.loads(response))
        st.table(df)

if uploaded_file is not None:
    # 读取文件内容
    query = read_file(uploaded_file)

    # 显示上传的文件内容
    st.chat_message("user").write(query)

    # 使用st.empty()创建一个可更新的区域
    empty_placeholder = st.empty()

    # 在该区域内显示“正在提取简历信息，请稍候...”
    with empty_placeholder.container():
        st.write("正在提取简历信息，请稍候...")

    try:
        # 准备prompt
        prompt = template.format(query=query).strip()
        prompt += "<sep>"
        inputs = tokenizer(prompt, return_tensors="pt")["input_ids"].cuda()
        max_new_tokens = 1024 - len(inputs[0])  # 确保总长度不超过1024
        outputs = model.generate(inputs, do_sample=False, max_new_tokens=max_new_tokens)  # 设置解码方式和最大生成长度
        output = tokenizer.decode(outputs[0])
        response = output.split("<sep>")[-1].replace("<eod>", '').strip()

        # 处理长度不同的列表
        parsed_response = json.loads(response)
        max_length = max(len(entities) for entities in parsed_response.values())
        for key in parsed_response:
            parsed_response[key] = parsed_response[key] + [''] * (max_length - len(parsed_response[key]))

        # 创建DataFrame
        df = pd.DataFrame(parsed_response)

        # 清空“正在提取简历信息，请稍候...”的信息
        empty_placeholder.empty()

        # 创建一个新的st.chat_message并显示解析后的表格
        new_placeholder = st.chat_message("assistant")
        with new_placeholder.container():
            st.table(df)
    except Exception as e:
        # 清空“正在提取简历信息，请稍候...”的信息
        empty_placeholder.empty()

        # 创建一个新的st.chat_message并显示错误信息
        new_placeholder = st.chat_message("assistant")
        with new_placeholder.container():
            st.error(f"发生了一个错误：{str(e)}")

# 如果没有上传文件，则显示提示信息
else:
    st.info("请上传一份简历文件以便开始分析。")